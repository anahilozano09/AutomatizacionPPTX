import pandas as pd
import hashlib
import os
import re
from pptx import Presentation
import win32com.client
from datetime import datetime
import time

# Diccionario para los meses en español
meses = {
    1: "Enero",
    2: "Febrero",
    3: "Marzo",
    4: "Abril",
    5: "Mayo",
    6: "Junio",
    7: "Julio",
    8: "Agosto",
    9: "Septiembre",
    10: "Octubre",
    11: "Noviembre",
    12: "Diciembre"
}


# Documentos a utilizar
archivoExcel = "DatosConstancias.xlsx"
plantillaEvento = "Constancia2026-2.pptx"
plantillaInter = "Intersem2026-2.pptx"
carpetaSalida = "Constancias 2026-2"

#Se crea la carpeta raíz o principal de las constancias
os.makedirs(carpetaSalida, exist_ok=True)

# Primera parte del proceso: generación de las funciones para nuestra automatizacion

# Función para generar firmas digitales
def generar_hash(fila):
    texto = f"{fila['Nombre']}_{fila['NombreEvento']}_{fila['Fecha']}"
    
    # Generar SHA-256
    return hashlib.sha256(texto.encode('utf-8')).hexdigest()

# Función para obtener la información de la fila que contiene la firma digital del documento Excel
def obtener_firma(fila):
    if pd.notna(fila['FirmaDigital']) and str(fila['FirmaDigital']).strip() != '':
        return fila['FirmaDigital']
    else:
        return generar_hash(fila)

# Función para evitar caracteres invalidos para la generación de archivos
def limpiar_nombre(nombre):
    return re.sub(r'[\\/*?:"<>|]', "", str(nombre))

# Función para reemplazar texto en el archivo de PowerPoint
def reemplazar_texto(shape, datos):
    if not shape.has_text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        full_text = "".join(run.text for run in paragraph.runs)

        for key, value in datos.items():
            placeholder = f"{{{{{key}}}}}"
            if placeholder in full_text:
                full_text = full_text.replace(placeholder, str(value))

        # Limpiar runs
        for run in paragraph.runs:
            run.text = ""

        # Escribir todo en el primer run
        if paragraph.runs:
            paragraph.runs[0].text = full_text

# Función para generar texto dependiendo el tipo de participación y evento
def generar_texto_participacion(fila):
    tipoPart = str(fila['TipoParticipacion']).strip().lower()
    tipoEvento = ""
    tipoEventoCol = str(fila['TipoEvento']).strip().lower()

    evento = str(fila['NombreEvento']).strip()


    if (tipoEventoCol == "curso") or (tipoEventoCol == "intersemestral"):
        tipoEvento = "el curso"
    elif tipoEventoCol == "conferencia":
        tipoEvento = "la conferencia"
    elif tipoEventoCol == "foro":
        tipoEvento = "el foro"
    elif tipoEventoCol == "evento":
        tipoEvento = "el evento"

    #Texto completo
    if tipoPart == "asistencia":
        return f"Por su asistencia en {tipoEvento} “{evento}”"
    
    elif tipoPart == "colaboracion":
        return f"Por su colaboración en {tipoEvento} “{evento}”"
    
    elif tipoPart == "participacion":
        if tipoEventoCol == "intersemestral":
            return f"Por haber en {tipoEvento} “{evento}”"
        else:
            return f"Por su participación en {tipoEvento} “{evento}”"
    
    elif tipoPart == "ponente":
        if (tipoEventoCol == "curso") or (tipoEventoCol == "conferencia") or (tipoEventoCol == "intersemestral"):
            return f"Por impartir {tipoEvento} “{evento}”"
        elif tipoEventoCol == "foro":
            return f"Por su participación como ponente en {tipoEvento} “{evento}”"
        else:
            return f"Por su participación como ponente en “{evento}”"
    else:
        return f"Por su participación en {tipoEvento} “{evento}”"

# Función para el formato de fecha de elaboracion para las constancias
def formato_fecha():
    hoy = datetime.now()
    dia = hoy.day
    mes = meses[hoy.month]
    anio = hoy.year

    return f"{dia} de {mes} del {anio}"

# Función para obtener la fecha de elaboración de las constancias (fecha en que se ejecuto el código)
def obtener_fecha_elaboracion(fila):
    if pd.notna(fila.get('FechaElaboracion')) and str(fila['FechaElaboracion']).strip() != "":
        return fila['FechaElaboracion']
    else:
        return formato_fecha()
    
# Función para cambiar el formato de la fecha del curso
def formatear_fecha_evento(fecha):
    try:
        # Si viene como NaN o vacío
        if pd.isna(fecha):
            return ""

        # Convertir a datetime si es string
        fecha = pd.to_datetime(fecha)

        dia = fecha.day
        mes = meses[fecha.month]
        anio = fecha.year

        return f"{dia} de {mes} del {anio}"
    
    except Exception:
        return str(fecha)
    
#Función para definir el rango de fechas de un inter
def formatear_rango_fechas(fecha_inicio, fecha_fin):
    try:
        f1 = pd.to_datetime(fecha_inicio)
        f2 = pd.to_datetime(fecha_fin)

        dia_inicio = f1.day
        dia_fin = f2.day
        mes_inicio = meses.get(f1.month, "")
        mes_fin = meses.get(f2.month, "")
        anio = f1.year

        if mes_inicio == mes_fin:
            return f"del {dia_inicio:02d} al {dia_fin:02d} de {mes_inicio} de {anio}"
        else:
            return f"del {dia_inicio:02d} del {mes_inicio} al {dia_fin:02d} del {mes_fin} de {anio}"
    except:
        return ""

# Texto para intersemestrales
def generar_texto_inter(fila):
    evento = str(fila['NombreEvento']).strip()
    semestre = str(fila['Semestre']).strip()

    rango_fechas = formatear_rango_fechas(fila['Fecha_dt'], fila['FechaFin_dt'])

    return f"Por haber aprobado el curso “{evento}” impartido {rango_fechas}, perteneciente a los cursos intersemestrales {semestre}"


# Función para generar nombres unicos para los archivos de las constancias
def generar_nombre_unico(ruta_base):
    if not os.path.exists(ruta_base):
        return ruta_base
    
    contador = 1
    nombre, ext = os.path.splitext(ruta_base)
    
    while True:
        nueva_ruta = f"{nombre}_{contador}{ext}"
        if not os.path.exists(nueva_ruta):
            return nueva_ruta
        contador += 1


#Segunda parte del proceso: Lectura, creación de las firmas digitales, modificación/actualización del archivo de Excel

df = pd.read_excel(archivoExcel)

if 'ArchivoPDF' not in df.columns:
    df['ArchivoPDF'] = ""

if 'FirmaDigital' not in df.columns:
    df['FirmaDigital'] = ""

#Generación de las firmas digitales
df['FirmaDigital'] = df.apply(obtener_firma, axis = 1)

if 'FechaElaboracion' not in df.columns:
    df['FechaElaboracion'] = ""

if 'TextoFinal' not in df.columns:
    df['TextoFinal'] = ""

#Copias de las fechas antes de formatear
df['Fecha_dt'] = pd.to_datetime(df['Fecha'], format="%d/%m/%Y", errors='coerce')
df['FechaFin_dt'] = pd.to_datetime(df['FechaFin'], format="%d/%m/%Y", errors='coerce')


#Formato para la fecha del evento
df['Fecha'] = df['Fecha'].apply(formatear_fecha_evento)

#Formato para la fecha fin del inter
df['FechaFin'] = df['FechaFin'].apply(formatear_fecha_evento)

#Obtencion de la fecha de elaboracion
df['FechaElaboracion'] = df.apply(obtener_fecha_elaboracion, axis = 1)

#Se coloca aquí el texto dependiendo de la participación y evento
df['TextoParticipacion'] = df.apply(generar_texto_participacion, axis = 1)

df['HorasComp'] = df['HorasComp'].apply(
    lambda x: str(int(x)) if pd.notna(x) and float(x).is_integer() else str(x)
)


#Texto final para el caso de los intersemestrales
df['TextoFinal'] = df.apply(
    lambda fila: generar_texto_inter(fila)
    if (str(fila['TipoEvento']).strip().lower() == "intersemestral") and (str(fila['TipoParticipacion']).strip().lower() == "participacion")
    else fila['TextoParticipacion'],
    axis = 1
)

#Actualizacion del Excel
df.to_excel(archivoExcel, index = False)
print("Excel actualizado")

#Tercera parte del proceso: Creación de carpetas por curso, generación de la constancia en pptx y guardado de las mismas

archivosGenerados = []

for index, fila in df.iterrows():
    datos = fila.to_dict()

    nombreLimpio = limpiar_nombre(fila['Nombre'])
    tipoEvento = str(fila['TipoEvento']).strip().lower()
    evento = limpiar_nombre(fila['NombreEvento'])
    fecha = limpiar_nombre(fila['Fecha'])
    hashCorto = fila['FirmaDigital'][:8]

    # Creacion de carpeta por curso
    carpetaEvento = os.path.join(carpetaSalida, evento)
    os.makedirs(carpetaEvento, exist_ok=True)

    #Nombre base del archivo
    archivoBase = os.path.join(
        carpetaEvento,
        f"{nombreLimpio}_{hashCorto}.pptx"
    )

    #Verificamos si ya existe el archivo
    if os.path.exists(archivoBase):
        print("La constancia ya existe")
        continue

    #Para evitar la sobreescritura de archivos
    archivoPower = generar_nombre_unico(archivoBase)
    

    #Creación de la constancia pptx

    if tipoEvento == "intersemestral":
        prs = Presentation(plantillaInter)
    
    else:
        prs = Presentation(plantillaEvento)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            reemplazar_texto(shape, datos)
        
    prs.save(archivoPower)
    archivosGenerados.append(archivoPower)

print("Constancias pptx generadas")

# Cuarta parte del proceso: Convertir las constancias pptx a pdf

powerpoint = win32com.client.Dispatch("PowerPoint.Application")
powerpoint.Visible = 1

for index, archivo in enumerate(archivosGenerados):
    fila = df.iloc[index]
    try:
        ruta_pptx = os.path.abspath(archivo)
        ruta_pdf = ruta_pptx.replace(".pptx", ".pdf")

        # Esperar un poco para asegurar que el archivo existe
        time.sleep(1)

        presentation = powerpoint.Presentations.Open(ruta_pptx, WithWindow=False)
        
        # 32 = formato PDF
        presentation.SaveAs(ruta_pdf, 32)
        presentation.Close()

        nombre_pdf = os.path.basename(ruta_pdf)
        df.loc[df['FirmaDigital'] == fila['FirmaDigital'], 'ArchivoPDF'] = nombre_pdf

        # Eliminar el archivo pptx
        if os.path.exists(ruta_pptx):
            os.remove(ruta_pptx)


    except Exception as e:
        print(f"Error con {archivo}: {e}")

powerpoint.Quit()

print("Conversión terminada")

df.to_excel(archivoExcel, index=False)
print("Excel actualizado con nombres de PDF")